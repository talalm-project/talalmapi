from types import SimpleNamespace

from app.operations.embeddings.generate_local_embeddings import GenerateLocalEmbeddings


class FakeLlama:
    calls = []

    def __init__(self, **kwargs):
        self.kwargs = kwargs

    def create_embedding(self, input):
        self.calls.append(input)
        return {"data": [{"embedding": [float(len(input)), 1.0]}]}


def test_generate_local_embeddings_returns_vector_ready_records_for_txt(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("First paragraph.\n\nSecond paragraph.", encoding="utf-8")
    model = SimpleNamespace(
        embedding_name="Qwen Embedding",
        embedding_local_file_path="/models/qwen-embedding.gguf",
    )
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model=model,
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert FakeLlama.calls == ["First paragraph.\n\nSecond paragraph."]
    assert operation.embeddings == [
        {
            "text": "First paragraph.\n\nSecond paragraph.",
            "embedding": [35.0, 1.0],
            "metadata": {
                "source": str(input_file),
                "source_name": "notes.txt",
                "extension": ".txt",
                "chunk_index": 0,
                "model": "Qwen Embedding",
                "model_path": "/models/qwen-embedding.gguf",
            },
        }
    ]


def test_generate_local_embeddings_removes_nul_bytes_from_text(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("First\x00 paragraph.", encoding="utf-8")
    model = SimpleNamespace(
        embedding_name="Qwen Embedding",
        embedding_local_file_path="/models/qwen-embedding.gguf",
    )
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model=model,
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert FakeLlama.calls == ["First paragraph."]
    assert operation.embeddings[0]["text"] == "First paragraph."


def test_generate_local_embeddings_accepts_manifest_model_dict(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("Content", encoding="utf-8")
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Manifest Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert operation.embeddings[0]["metadata"]["model"] == "Manifest Embedding"
    assert operation.embeddings[0]["metadata"]["model_path"] == "/models/embedding.gguf"


def test_generate_local_embeddings_defaults_model_name_from_model_path(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("Content", encoding="utf-8")
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"path": "/models/qwen-embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert operation.embeddings[0]["metadata"]["model"] == "qwen-embedding"


def test_generate_local_embeddings_extracts_docx(monkeypatch, tmp_path):
    from docx import Document

    input_file = tmp_path / "notes.docx"
    document = Document()
    document.add_paragraph("Document paragraph")
    document.save(input_file)
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert FakeLlama.calls == ["Document paragraph"]
    assert operation.embeddings[0]["metadata"]["paragraph"] == 1


def test_generate_local_embeddings_extracts_pptx(monkeypatch, tmp_path):
    from pptx import Presentation

    input_file = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slide title"
    presentation.save(input_file)
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert FakeLlama.calls == ["Slide title"]
    assert operation.embeddings[0]["metadata"]["slide"] == 1


def test_generate_local_embeddings_extracts_xlsx(monkeypatch, tmp_path):
    from openpyxl import Workbook

    input_file = tmp_path / "sheet.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Facts"
    sheet["A1"] = "Cell"
    sheet["B1"] = "Value"
    workbook.save(input_file)
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert FakeLlama.calls == ["Cell\tValue"]
    assert operation.embeddings[0]["metadata"]["sheet"] == "Facts"


def test_generate_local_embeddings_chunks_large_text(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("abcdefghij", encoding="utf-8")
    FakeLlama.calls = []
    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FakeLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
        chunk_size=4,
        chunk_overlap=1,
    )
    operation.execute()

    assert operation.valid()
    assert FakeLlama.calls == ["abcd", "defg", "ghij"]
    assert [record["metadata"]["chunk_index"] for record in operation.embeddings] == [0, 1, 2]


def test_generate_local_embeddings_splits_chunks_to_llama_token_limit(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("abcdefghij", encoding="utf-8")
    calls = []

    class TokenLimitedLlama:
        n_batch = 5

        def __init__(self, **kwargs):
            pass

        def n_ctx(self):
            return 5

        def tokenize(self, text, add_bos=True, special=False):
            tokens = list(text)
            return ([0] if add_bos else []) + tokens

        def detokenize(self, tokens, prev_tokens=None, special=False):
            return bytes(tokens)

        def create_embedding(self, input):
            calls.append(input)
            return {"data": [{"embedding": [float(len(input))]}]}

    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: TokenLimitedLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert calls == ["abcd", "efgh", "ij"]
    assert [record["metadata"]["chunk_index"] for record in operation.embeddings] == [0, 1, 2]


def test_generate_local_embeddings_uses_llama_embedding_sequence_slot_limit(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("abcdefghij", encoding="utf-8")
    calls = []

    class ContextParams:
        n_ctx = 12
        n_seq_max = 3

    class SlotLimitedLlama:
        n_batch = 12
        context_params = ContextParams()

        def __init__(self, **kwargs):
            pass

        def tokenize(self, text, add_bos=True, special=False):
            tokens = list(text)
            return ([0] if add_bos else []) + tokens

        def detokenize(self, tokens, prev_tokens=None, special=False):
            return bytes(tokens)

        def create_embedding(self, input):
            calls.append(input)
            return {"data": [{"embedding": [float(len(input))]}]}

    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: SlotLimitedLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert calls == ["abc", "def", "ghi", "j"]


def test_generate_local_embeddings_prefers_runtime_context_for_sequence_slot_limit(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("abcdefghij", encoding="utf-8")
    calls = []

    class ContextParams:
        n_ctx = 4
        n_seq_max = 4

    class RuntimeContextLlama:
        n_batch = 8
        context_params = ContextParams()

        def __init__(self, **kwargs):
            pass

        def n_ctx(self):
            return 16

        def tokenize(self, text, add_bos=True, special=False):
            tokens = list(text)
            return ([0] if add_bos else []) + tokens

        def detokenize(self, tokens, prev_tokens=None, special=False):
            return bytes(tokens)

        def create_embedding(self, input):
            calls.append(input)
            return {"data": [{"embedding": [float(len(input))]}]}

    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: RuntimeContextLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.valid()
    assert calls == ["abc", "def", "ghi", "j"]


def test_generate_local_embeddings_uses_configured_max_input_tokens(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("abcdefghij", encoding="utf-8")
    calls = []

    class WideContextLlama:
        n_batch = 12

        def __init__(self, **kwargs):
            pass

        def n_ctx(self):
            return 12

        def tokenize(self, text, add_bos=True, special=False):
            tokens = list(text)
            return ([0] if add_bos else []) + tokens

        def detokenize(self, tokens, prev_tokens=None, special=False):
            return bytes(tokens)

        def create_embedding(self, input):
            calls.append(input)
            return {"data": [{"embedding": [float(len(input))]}]}

    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: WideContextLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
        max_input_tokens=4,
    )
    operation.execute()

    assert operation.valid()
    assert calls == ["abc", "def", "ghi", "j"]


def test_generate_local_embeddings_returns_error_when_llama_decode_fails(monkeypatch, tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("Content", encoding="utf-8")

    class FailingLlama:
        def __init__(self, **kwargs):
            pass

        def create_embedding(self, input):
            raise RuntimeError("llama_decode returned 1")

    monkeypatch.setattr("app.operations.embeddings.generate_local_embeddings._llama_class", lambda: FailingLlama)

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.invalid()
    assert operation.embeddings == []
    assert operation.errors["embedding"] == ["llama-cpp failed: llama_decode returned 1"]


def test_generate_local_embeddings_rejects_missing_model_path(tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("Content", encoding="utf-8")

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.invalid()
    assert operation.errors["embedding_local_file_path"] == ["required"]


def test_generate_local_embeddings_rejects_non_gguf_model_path(tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("Content", encoding="utf-8")

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.bin"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.invalid()
    assert operation.errors["embedding_local_file_path"] == ["must be a .gguf model"]


def test_generate_local_embeddings_rejects_unsupported_input_file(tmp_path):
    input_file = tmp_path / "notes.csv"
    input_file.write_text("Content", encoding="utf-8")

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
    )
    operation.execute()

    assert operation.invalid()
    assert operation.errors["input_file"] == ["unsupported file type"]


def test_generate_local_embeddings_rejects_invalid_chunk_options(tmp_path):
    input_file = tmp_path / "notes.txt"
    input_file.write_text("Content", encoding="utf-8")

    operation = GenerateLocalEmbeddings(
        local_embedding_model={"name": "Embedding", "path": "/models/embedding.gguf"},
        input_file=input_file,
        chunk_size=10,
        chunk_overlap=10,
    )
    operation.execute()

    assert operation.invalid()
    assert operation.errors["chunk_overlap"] == ["must be less than chunk_size"]
