from pathlib import Path

from app.operations.connectors.metadata import embedding_local_file_path, embedding_model_name


SUPPORTED_INPUT_EXTENSIONS = {".txt", ".pptx", ".docx", ".xlsx", ".pdf"}
DEFAULT_CHUNK_SIZE = 1200
DEFAULT_CHUNK_OVERLAP = 120


class MissingParserDependency(Exception):
    def __init__(self, package_name):
        super().__init__(package_name)
        self.package_name = package_name


class GenerateLocalEmbeddings:
    def __init__(
        self,
        local_embedding_model,
        input_file,
        chunk_size=DEFAULT_CHUNK_SIZE,
        chunk_overlap=DEFAULT_CHUNK_OVERLAP,
        model_options=None,
        max_input_tokens=None,
        source_name=None,
    ):
        self.local_embedding_model = local_embedding_model
        self.input_file = input_file
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.model_options = model_options or {}
        self.max_input_tokens = max_input_tokens
        self.source_name = source_name
        self.errors = {}
        self.chunks = []
        self.embeddings = []

    def execute(self):
        self.errors = self._validation_errors()
        self.chunks = []
        self.embeddings = []
        if self.errors:
            return

        try:
            blocks = self._extract_text_blocks()
        except MissingParserDependency as error:
            self.errors = {"input_file": [f"missing parser dependency: {error.package_name}"]}
            return

        self.chunks = self._chunk_blocks(blocks)
        if not self.chunks:
            self.errors = {"input_file": ["no text found"]}
            return

        try:
            llm = _llama_class()(model_path=self._model_path(), embedding=True, **self.model_options)
            self.chunks = _split_chunks_for_llama(self.chunks, llm, self.max_input_tokens)
            if not self.chunks:
                self.errors = {"input_file": ["no text found"]}
                return

            self.embeddings = [
                {
                    "text": chunk["text"],
                    "embedding": _embedding_from_response(llm.create_embedding(chunk["text"])),
                    "metadata": self._metadata_for_chunk(chunk, index),
                }
                for index, chunk in enumerate(self.chunks)
            ]
        except RuntimeError as error:
            self.embeddings = []
            self.errors = {"embedding": [f"llama-cpp failed: {error}"]}

    def valid(self):
        return not self.errors

    def invalid(self):
        return bool(self.errors)

    def _validation_errors(self):
        errors = {}
        model_path = self._model_path()
        input_path = self._input_path()

        if not model_path:
            errors["embedding_local_file_path"] = ["required"]
        elif not model_path.lower().endswith(".gguf"):
            errors["embedding_local_file_path"] = ["must be a .gguf model"]

        if input_path is None:
            errors["input_file"] = ["required"]
        elif not input_path.exists():
            errors["input_file"] = ["not found"]
        elif not input_path.is_file():
            errors["input_file"] = ["invalid"]
        elif input_path.suffix.lower() not in SUPPORTED_INPUT_EXTENSIONS:
            errors["input_file"] = ["unsupported file type"]

        if not isinstance(self.chunk_size, int) or self.chunk_size <= 0:
            errors["chunk_size"] = ["must be greater than 0"]
        if not isinstance(self.chunk_overlap, int) or self.chunk_overlap < 0:
            errors["chunk_overlap"] = ["must be greater than or equal to 0"]
        if (
            isinstance(self.chunk_size, int)
            and isinstance(self.chunk_overlap, int)
            and self.chunk_size > 0
            and self.chunk_overlap >= self.chunk_size
        ):
            errors["chunk_overlap"] = ["must be less than chunk_size"]
        if not isinstance(self.model_options, dict):
            errors["model_options"] = ["invalid"]
        if self.max_input_tokens is not None and (not isinstance(self.max_input_tokens, int) or self.max_input_tokens <= 0):
            errors["max_input_tokens"] = ["must be greater than 0"]

        return errors

    def _extract_text_blocks(self):
        input_path = self._input_path()
        extension = input_path.suffix.lower()

        if extension == ".txt":
            return _extract_txt(input_path)
        if extension == ".docx":
            return _extract_docx(input_path)
        if extension == ".pptx":
            return _extract_pptx(input_path)
        if extension == ".xlsx":
            return _extract_xlsx(input_path)
        if extension == ".pdf":
            return _extract_pdf(input_path)

        return []

    def _chunk_blocks(self, blocks):
        chunks = []
        for block in blocks:
            text = _clean_text(block.get("text"))
            if not text:
                continue

            chunks.extend(
                {
                    "text": chunk_text,
                    "metadata": block.get("metadata", {}),
                }
                for chunk_text in _chunk_text(text, self.chunk_size, self.chunk_overlap)
            )

        return chunks

    def _metadata_for_chunk(self, chunk, index):
        input_path = self._input_path()
        metadata = {
            "source": str(input_path),
            "source_name": self.source_name or input_path.name,
            "extension": input_path.suffix.lower(),
            "chunk_index": index,
            "model": self._model_name(),
            "model_path": self._model_path(),
        }
        metadata.update(chunk.get("metadata", {}))
        return metadata

    def _input_path(self):
        if self.input_file is None:
            return None

        try:
            return Path(self.input_file)
        except TypeError:
            return None

    def _model_path(self):
        metadata_path = embedding_local_file_path(self.local_embedding_model)
        if metadata_path:
            return metadata_path

        return _first_present(self.local_embedding_model, ["local_file_path", "path", "model_path"])

    def _model_name(self):
        name = embedding_model_name(self.local_embedding_model) or _first_present(self.local_embedding_model, ["name", "model"])
        if name:
            return name

        model_path = self._model_path()
        if model_path:
            return Path(model_path).stem

        return None


def _extract_txt(input_path):
    return [{"text": input_path.read_text(encoding="utf-8"), "metadata": {}}]


def _extract_docx(input_path):
    try:
        from docx import Document
    except ImportError as error:
        raise MissingParserDependency("python-docx") from error

    document = Document(input_path)
    blocks = []
    for index, paragraph in enumerate(document.paragraphs):
        text = _clean_text(paragraph.text)
        if text:
            blocks.append({"text": text, "metadata": {"paragraph": index + 1}})

    for table_index, table in enumerate(document.tables):
        rows = []
        for row in table.rows:
            values = [_clean_text(cell.text) for cell in row.cells]
            values = [value for value in values if value]
            if values:
                rows.append("\t".join(values))
        if rows:
            blocks.append({"text": "\n".join(rows), "metadata": {"table": table_index + 1}})

    return blocks


def _extract_pptx(input_path):
    try:
        from pptx import Presentation
    except ImportError as error:
        raise MissingParserDependency("python-pptx") from error

    presentation = Presentation(input_path)
    blocks = []
    for slide_index, slide in enumerate(presentation.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = _clean_text(shape.text)
                if text:
                    texts.append(text)
        if texts:
            blocks.append({"text": "\n".join(texts), "metadata": {"slide": slide_index + 1}})

    return blocks


def _extract_xlsx(input_path):
    try:
        from openpyxl import load_workbook
    except ImportError as error:
        raise MissingParserDependency("openpyxl") from error

    workbook = load_workbook(input_path, read_only=True, data_only=True)
    blocks = []
    try:
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [_clean_text(value) for value in row]
                values = [value for value in values if value]
                if values:
                    rows.append("\t".join(values))
            if rows:
                blocks.append({"text": "\n".join(rows), "metadata": {"sheet": sheet.title}})
    finally:
        workbook.close()

    return blocks


def _extract_pdf(input_path):
    try:
        from pypdf import PdfReader
    except ImportError as error:
        raise MissingParserDependency("pypdf") from error

    reader = PdfReader(input_path)
    blocks = []
    for page_index, page in enumerate(reader.pages):
        text = _clean_text(page.extract_text())
        if text:
            blocks.append({"text": text, "metadata": {"page": page_index + 1}})

    return blocks


def _chunk_text(text, chunk_size, chunk_overlap):
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - chunk_overlap

    return chunks


def _split_chunks_for_llama(chunks, llm, max_input_tokens=None):
    if not hasattr(llm, "tokenize") or not hasattr(llm, "detokenize"):
        return chunks

    runtime_max_tokens = _llama_max_input_tokens(llm)
    if max_input_tokens and runtime_max_tokens:
        max_tokens = min(max_input_tokens, runtime_max_tokens)
    else:
        max_tokens = max_input_tokens or runtime_max_tokens
    if max_tokens is None:
        return chunks

    max_content_tokens = max_tokens - 1
    if max_content_tokens < 1:
        return chunks

    split_chunks = []
    for chunk in chunks:
        tokens = llm.tokenize(chunk["text"].encode("utf-8"), add_bos=False)
        if len(tokens) + 1 <= max_tokens:
            split_chunks.append(chunk)
            continue

        for start in range(0, len(tokens), max_content_tokens):
            text = llm.detokenize(tokens[start : start + max_content_tokens]).decode("utf-8", errors="replace").strip()
            if text:
                split_chunks.append({"text": text, "metadata": chunk.get("metadata", {})})

    return split_chunks


def _llama_max_input_tokens(llm):
    limits = []
    n_batch = getattr(llm, "n_batch", None)
    if isinstance(n_batch, int) and n_batch > 0:
        limits.append(n_batch)

    runtime_n_ctx = None
    n_ctx_method = getattr(llm, "n_ctx", None)
    if callable(n_ctx_method):
        value = n_ctx_method()
        if isinstance(value, int) and value > 0:
            runtime_n_ctx = value

    context_params = getattr(llm, "context_params", None)
    n_seq_max = getattr(context_params, "n_seq_max", None)
    if runtime_n_ctx:
        limits.append(runtime_n_ctx)
        if isinstance(n_seq_max, int) and n_seq_max > 0:
            limits.append(max(runtime_n_ctx // n_seq_max, 1))
    else:
        n_ctx = getattr(context_params, "n_ctx", None)
        if isinstance(n_ctx, int) and n_ctx > 0:
            limits.append(n_ctx)
            if isinstance(n_seq_max, int) and n_seq_max > 0:
                limits.append(max(n_ctx // n_seq_max, 1))

    if not limits:
        return None

    return min(limits)


def _embedding_from_response(response):
    data = _get_value(response, "data")
    if not data:
        return None

    first_record = data[0]
    return _get_value(first_record, "embedding")


def _first_present(source, keys):
    for key in keys:
        value = _get_value(source, key)
        if value:
            return value

    return None


def _get_value(source, key):
    if isinstance(source, dict):
        return source.get(key)

    return getattr(source, key, None)


def _clean_text(value):
    if value is None:
        return ""

    return str(value).replace("\x00", "").strip()


def _llama_class():
    from llama_cpp import Llama

    return Llama
